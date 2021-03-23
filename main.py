from flask import Flask,render_template,request,session,redirect,url_for,jsonify,make_response,send_from_directory
import werkzeug
import os
import sys
import codecs
from datetime import datetime
# PP
from pptx import Presentation
from pptx.util import Inches
# 各処理の取得
# 男女
import Q1_man_female
# １列テキスト
import Q11_single_text
# 2列テキスト
import Q12_twice_text
# １列掲載許可あり
import Qextra_nostalgicspeed_singletext
# 2列掲載許可あり
import Qextra_nostalgicspeed_doubletext
# ラジオボタン＋テキスト回答
import Q9_two_one_select
# チェックボックス
import Q5_bestpages
# マツダファンブック固有
import Qextra_matsudafanbook
# カミオン固有
import Q3_camion_prefectures
# テスト
app = Flask(__name__)

UPLOAD_DIR = "./samples"
ALLOWED_EXTENTIONS = set(['csv','xlsx'])
# ファイルサイズ：1MB
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024
app.config['UPLOAD_FOLDER'] = UPLOAD_DIR

@app.route('/')
def top():
    try:
        return render_template('index.html')
    except Exception as e:
        return str(e)

@app.route('/data_print',methods=['POST'])
def data_print():
    if request.method == 'POST':
        # 選択したファイルの取得、保存
        get_files = request.files.getlist('uploads')
        for file in get_files:
            fileName = file.filename
            saveFileName = werkzeug.utils.secure_filename(fileName)
            file.save(os.path.join(UPLOAD_DIR, saveFileName))

        # template(PP)の初期化（一枚目のテンプレート以外を削除）
        prs = Presentation('./static/pp/template.pptx')
        if len(prs.slides) > 1:
            slide = prs.slides
            counts = 0
            for s in slide:
                if counts > 0:
                    id_dict = {slide.id: [i, slide.rId] for i,slide in enumerate(prs.slides._sldIdLst)}
                    slide_id = s.slide_id
                    prs.part.drop_rel(id_dict[slide_id][1])
                    del prs.slides._sldIdLst[id_dict[slide_id][0]]
                counts+=1
            # 削除後、削除した状態でPPを保存
            prs.save('./static/pp/template.pptx')
        
        # 設問番号、設問内容のリストを取得
        req_num = request.form.getlist('q_num')
        req_sel = request.form.getlist('q_select')
        for nums,sels in zip(req_num,req_sel):
            if sels == 'radio':
                Q1_man_female.male_female(nums)
            elif sels == 'singletext':
                Q11_single_text.singletext(nums)
            elif sels == 'doubletext':
                Q12_twice_text.doubletext(nums)
            elif sels == 'permit_singletext':
                Qextra_nostalgicspeed_singletext.permit_single(nums)
            elif sels == 'permit_doubletext':
                Qextra_nostalgicspeed_doubletext.permit_double(nums)
            elif sels == 'radio_plus_text':
                Q9_two_one_select.graph_chart_text(nums)
            elif sels == 'checkbox':
                Q5_bestpages.checkbox(nums)
            elif sels == 'matsuda':
                Qextra_matsudafanbook.Matsuda(nums)
            elif sels == 'camion':
                Q3_camion_prefectures.Camion(nums)
        return render_template('data_print.html')
if __name__ == '__main__':
    app.run(host="0.0.0.0",debug=True,port=int(os.environ.get("PORT",5000)))