# 職業（回答パターン：３[通常、付随テキスト、その他]）
import csv
import codecs
import locale
#グラフ用
import matplotlib
import matplotlib.pyplot as plt
import japanize_matplotlib
# データの扱いに必要
import pandas as pd
import numpy as np
import datetime as dt
# 画像の保存先の指定に必要
import os
# PP（パワーポイント）の作成、挿入
from pptx import Presentation
from pptx.util import Inches
# 画像の読み込み
from glob import glob
# グラフ用
import seaborn as sns
import re
# 分割した機能
import Q4_parts.Q4_main_chart as Q4_main_chart
import Q4_parts.Q4_chart_twice as Q4_chart_twice
import Q4_parts.Q4_chart_third as Q4_chart_third
import Q4_parts.Q4_graph as Q4_graph

# グラフの作成、実行
Q4_main_chart.Q4_main_chart()
Q4_chart_twice.Q4_chart_twice()
Q4_chart_third.Q4_chart_third()
Q4_graph.Q4_graph()

# PPを開いて画像を張り込む
# 作成した画像をPPに貼り付ける
# PPの呼び出し
prs = Presentation('./template.pptx')
# 画像を中央配置にするためにスライドの幅と高さを取得する
width = prs.slide_width
height = prs.slide_height
# 貼り付ける画像の読み込み
fnms = ['./save_images/Q4_main_chart.png','./save_images/Q4_chart_two.png','./save_images/Q4_chart_third.png','./save_images/Q4_graph.png']
# ファイルでループさせる
for fnm in fnms:
    # スライドの追加。６は白紙のスライド
    add_s = prs.slide_layouts[6]
    # 合体？
    slide = prs.slides.add_slide(add_s)
    # 画像の挿入
    pic = slide.shapes.add_picture(fnm,0,0,width=None,height=None)
    # 画像を中心に設置
    pic.left = int((width - pic.width)/2)
    pic.top = int((height - pic.height)/2)
# とりあえず画像サイズについては個々人で調整してもらう方向で
prs.save('./template.pptx')