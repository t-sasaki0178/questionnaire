# 1枠テキスト掲載許可あり（ノススピ参照）
import csv
import codecs
#グラフ用
import matplotlib
import matplotlib.pyplot as plt
# データの扱いに必要
import pandas as pd
# 画像の保存先の指定に必要
import os
# PP（パワーポイント）の作成、挿入
from pptx import Presentation
from pptx.util import Inches
# 画像の読み込み
from glob import glob
import re

from jinja2 import Environment,FileSystemLoader 
from selenium import webdriver
from selenium.webdriver.chrome.options import Options
# driverまでのパスを求められることがあるが、以下１文のみで解決できる（理由は不明。おそらくモジュール内のdriverを使ってる？）
import chromedriver_binary
# nanの判定のため
import math
# listの順序を固定
from collections import OrderedDict

# フォントの指定（日本語が文字化けしないもの）
matplotlib.rcParams['font.family'] = 'IPAexGothic'

def permit_single(select):
    # csv読み込み版
    layout_load = pd.read_csv('./samples/layout.csv',encoding='cp932',header=0)
    # ldata:id（インデックス）、残り全部itemに入ってる
    # 最初に両方に値が入っているパターンを想定
    # 入っていない場合は入っている方を問題タイトルに、入っている場合は先に出てきたほうをタイトル、あとの方をタイトルのラベルにする
    l_titlelabel = []
    for ldata,item in layout_load.iterrows():
        if ('Q{}'.format(select) in item[1]) == True:
            if item[4] and item[5]:
                l_title = item[4]
                l_titlelabel.insert(len(l_titlelabel),item[5])
            else:
                if math.isnan(item[4]):
                    l_title = item[5]
                else:
                    l_title = item[4]

    # rawdataから問題の回答を取得する
    raw_load = pd.read_csv('./samples/rawdata.csv',header=None,encoding='shift_jisx0213')

    # rawdataの一列目から関連するインデックスを取得する
    # Q11に関連する部分を割り出す
    # 各idごとに最後の値（掲載許可）を取得
    rels = []
    ans = []
    # python3.7以降では仕様化されたはずのlistの順番が入れ替わるのでOrderedDict()で固定にする
    permits = OrderedDict()
    rcount = 0
    for rdata,vals in raw_load.iterrows():
        if rcount == 0:
            for rrans in vals:
                if ('Q{}'.format(select) in rrans) == True:
                    rels.append(1)
                else:
                    rels.append(0)
        else:
            ans.append(vals)
            try:
                if math.isnan(vals[len(vals)-1]):
                    permits[vals[0]] = '0'
            except:
                permits[vals[0]] = '1'
        rcount+=1

    # リストrelsにより関連する行のあぶり出しができたのでこのインデックス数を取得する
    index_num = [n for n,v in enumerate(rels) if v == 1]

    # laをforループさせたときに分離させたans（純回答データ）から
    # index_numを当てはめた引数変数を取得し、Q11の回答のみのリストを作成する
    # ２重の辞書型にして、idあたりのQ11に関連した値だけを{ID：辞書型}の形で入れる
    this_ans = {}
    ids_count = 1
    dict_count = 0
    for this in ans:
        for nums in index_num:
            if dict_count == 0:
                this_ans[ids_count] = {}
            try:
                if math.isnan(this[nums]):
                    pass
            except:
                this_ans[ids_count][nums] = this[nums]
            dict_count+=1
        ids_count+=1
        dict_count = 0

    # jinja2テンプレート出力
    env = Environment(loader=FileSystemLoader('./', encoding='utf8'))
    tmpl = env.get_template('jinja2_templetes/templetes_extra_nostalgicspeed.tmpl')
    # 商品情報を入れるリストを宣言、初期化
    items = []
    per_count = 0
    # ループでitemsにIDとテキストをappendする
    for tak,tav in this_ans.items():
        try:
            if math.isnan(tav):
                pass
        except:
            if len(tav) > 1:
                for val_key,val_value in tav.items():
                    permit = [int(per_val) for per_key,per_val in dict(permits).items() if str(tak) == per_key]
                    items.append({'ID':tak,'VALUE':val_value,'PERMIT':list(permit)[0]})
                per_count+=1
            else:
                for tav_keys,tav_values in tav.items():
                    permit = [int(per_val) for per_key,per_val in dict(permits).items() if str(tak) == per_key]
                    items.append({'ID':tak,'VALUE':tav_values,'PERMIT':list(permit)[0]})
                    per_count+=1


    #ここでtmplの変数に入る文字を指定している
    html = tmpl.render({'title':l_title,'item':items})
    with open('jinja2_templetes/Qextra_nostalgicspeed.html',mode='w') as f:
        f.write(str(html))

    # seleniumでブラウザ表示（動作問題なし）
    options = Options()
    options.add_argument('--headless')
    driver = webdriver.Chrome(options=options)
    driver.get('file:///Users/t_sasaki/Documents/%E6%A5%AD%E5%8B%99/%E3%82%A2%E3%83%B3%E3%82%B1%E3%83%BC%E3%83%88/%E3%82%AB%E3%83%9F%E3%82%AA%E3%83%B3/macro_mk/jinja2_templetes/Qextra_nostalgicspeed.html')
    page_width = driver.execute_script('return document.body.scrollWidth')
    page_height = driver.execute_script('return document.body.scrollHeight')
    driver.set_window_size(page_width,page_height)
    driver.save_screenshot('save_images/Q{}.png'.format(select))
    driver.quit()

    # PPを開いて画像を張り込む
    # 作成した画像をPPに貼り付ける
    # PPの呼び出し
    prs = Presentation('./static/pp/template.pptx')
    # 画像を中央配置にするためにスライドの幅と高さを取得する
    width = prs.slide_width
    height = prs.slide_height
    # 貼り付ける画像の読み込み
    fnms = ['./save_images/Q{}.png'.format(select)]
    # ファイルでループさせる
    for fnm in fnms:
        # スライドの追加。６は白紙のスライド
        add_s = prs.slide_layouts[6]
        # 合体？
        slide = prs.slides.add_slide(add_s)
        # 画像の挿入
        pic = slide.shapes.add_picture(fnm,0,0,width=None,height=None)
        # 画像を中心に設置
        pic.left = int((width - pic.width)/4)
        pic.top = int((height - pic.height)/4)
    # とりあえず画像サイズについては個々人で調整してもらう方向で
    prs.save('./static/pp/template.pptx')

if __name__ in '__main__':
    permit_single()