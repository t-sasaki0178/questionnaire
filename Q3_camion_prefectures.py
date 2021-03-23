# 都道府県択一
# カミオンバージョン（地方ごとに分ける、通常のグラフ＋表作成に固有の処理を追加）
import csv
import codecs
import locale
#グラフ用
import matplotlib
import matplotlib.pyplot as plt
import japanize_matplotlib
# データの扱いに必要
import pandas as pd
import numpy as np
import datetime as dt
# 画像の保存先の指定に必要
import os
# PP（パワーポイント）の作成、挿入
from pptx import Presentation
from pptx.util import Inches
# 画像の読み込み
from glob import glob
# グラフ用
import seaborn as sns
# 常用関数ファイル
from versatility import truncate

from jinja2 import Environment,FileSystemLoader 
from selenium import webdriver
from selenium.webdriver.chrome.options import Options
# driverまでのパスを求められることがあるが、以下１文のみで解決できる（理由は不明。おそらくモジュール内のdriverを使ってる？）
import chromedriver_binary
# nanの判定のため
import math

# seaborn:グラフとかがきれいにできる
sns.set()
sns.set_style('whitegrid')

# フォントの指定（日本語が文字化けしないもの）
matplotlib.rcParams['font.family'] = 'IPAexGothic'
plt.rcParams["font.size"] = 14

def Camion(select):
    # layoutデータから問題の回答を取得する
    layout_load = open('./samples/layout.csv',mode="r",encoding="cp932",errors='ignore',newline="")
    ll = csv.reader(layout_load)
    ll_outhead = next(ll)
    # 回答一覧のリストを保存用のリストに格納する
    # 表上かグラフ上に表示する問題テキストも取得しておく
    # save_lltextは文章をカットするのでカットされてない状態の文章も保存しておく
    save_lltext = []
    save_nocut_lltext = []
    for lltext in ll:
        if lltext[1] == 'Q{}'.format(select):
            save_lltext.append(lltext)
            save_nocut_lltext.append(lltext)
            if len(lltext[4]) == 0:
                graph_title = lltext[5]
            else:
                graph_title = lltext[4]

    # カットされてない選択肢テキスト。truncate()関数を先に動かすとこっちにも影響が出るので
    # truncate()関数の処理が行われる前にリストに格納する
    nocut_text = []
    th_counter = 0
    for snll in save_nocut_lltext:
        for snll_text in snll:
            if th_counter > 5:
                nocut_text.append(snll_text)
            th_counter+=1

    # そのままだと文字数が多いので、一定値以上の場合は文字数を短縮する
    for tl_index in range(0,len(save_lltext[0])):
        save_lltext[0][tl_index] = truncate(save_lltext[0][tl_index],10)

    # リストから選択肢がkeyになるように辞書型配列を作成する（値はカウントするので０で統一）
    # 回答部分を判断する変数とループ回数を判断する変数を用意
    # 値の変更（個別の件数）集計時にkey（選択肢）がないと変更できないので、keyのみの配列も作成する
    select_data = {}
    key_data = []
    through_counter = 0
    anser_counter = 1
    for sll in save_lltext:
        for sll_data in sll:
            if through_counter > 5:
                select_data[sll_data] = 0
                key_data.append(sll_data)
                anser_counter+=1
            through_counter+=1

    csv_load = open('./samples/rawdata.csv',mode="r",encoding="cp932",errors='ignore',newline="")
    f = csv.reader(csv_load)
    # headerを飛ばす（next関数はnextで飛ばした部分を元データから切り取る）
    f_outhead = next(f)

    # f_outhead（rawdataの１行目）から今回の問題に関連する部分（列番号）を抜き出す
    rel_number = 0
    rel_numcounter = 0
    for reldata in f_outhead:
        if reldata == 'Q{}'.format(select):
            rel_number = rel_numcounter
        rel_numcounter+=1

    # 取得した関連データの位置からそれぞれの回答を取得、select_dataのvalueへ加算させるために、key形式の値で直接入れる
    act_data = []
    for act in f:
        act_data.append(key_data[int(act[rel_number])-1])

    # 値ベースでループを回して、keyが出てくる度にvalueを加算していく
    # 回答データの配列act_data、関連した行数を格納したrel_number、key（選択肢）を格納したkey_data、辞書型配列（値はすべて初期値０）
    for adata in act_data:
        for s_key,s_value in select_data.items():
            if adata == s_key:
                select_data[adata]+=1
    # 回答に合わせた横棒グラフを作成
    # x:グラフの位置（横棒グラフの場合、数字が大きいほうが上に配置される）
    # range()はA<=range<Bなので初期値を1にして、ベースの値に1を加える
    x = []
    for ranges in range(1,len(key_data) + 1):
        x.append(ranges)
    # ラベルはリストの後ろから貼っていくので、key_dataを反転して挿入する
    x_label = []
    for rk in reversed(key_data):
        x_label.append(rk)
    # 値データの挿入。挿入順はx_label（逆順）に合わせる
    y = []
    # select_dataのvalueの値を逆順にして、yのlistに挿入する
    for s_ins in reversed(select_data.values()):
        y.append(s_ins)

    # x,x_label,yの値を地域ごとに再分割する
    # 北海道地方：北海道
    hokkaido = 0
    # 東北地方：青森県、岩手県、秋田県、宮城県、福島県、山形県
    tohoku = 0
    # 北関東地方：茨城県、栃木県、群馬県
    n_kanto = 0
    # 南関東地方：東京都、神奈川県、千葉県、埼玉県
    s_kanto = 0
    # 北陸地方：新潟県、富山県、石川県、福井県
    hokuriku = 0
    # 中央高地：山梨県、長野県
    center_hill = 0
    # 東海地方：岐阜県、静岡県、愛知県、三重県
    tokai = 0
    # 近畿地方：滋賀県、京都府、大阪府、兵庫県、奈良県、和歌山県
    kinki = 0
    # 中国地方：鳥取県、島根県、岡山県、広島県
    chugoku = 0
    # 四国地方：徳島県、香川県、愛媛県、高知県
    sikoku = 0
    # 九州地方：山口県、福岡県、佐賀県、長崎県、熊本県、大分県、宮崎県、鹿児島県、沖縄県
    kyusyu = 0

    district_label = ['九州地方','四国地方','中国地方','近畿地方','東海地方','中央高地','北陸地方','南関東地方','北関東地方','東北地方','北海道地方']
    for labels,numbers in zip(x_label,y):
        if labels == '北海道':
            hokkaido += int(numbers)
        elif labels == '青森県' or labels == '岩手県' or labels == '秋田県' or labels == '宮城県' or labels == '福島県' or labels == '山形県':
            tohoku += int(numbers)
        elif labels == '茨城県' or labels == '栃木県' or labels == '群馬県':
            n_kanto += int(numbers)
        elif labels == '東京都' or labels == '神奈川県' or labels == '千葉県' or labels == '埼玉県':
            s_kanto += int(numbers)
        elif labels == '新潟県' or labels == '富山県' or labels == '石川県' or labels == '福井県':
            hokuriku += int(numbers)
        elif labels == '山梨県' or labels == '長野県':
            center_hill += int(numbers)
        elif labels == '岐阜県' or labels == '静岡県' or labels == '愛知県' or labels == '三重県':
            tokai += int(numbers)
        elif labels == '滋賀県' or labels == '京都府' or labels == '大阪府' or labels == '兵庫県' or labels == '奈良県' or labels == '和歌山県':
            kinki += int(numbers)
        elif labels == '鳥取県' or labels == '島根県' or labels == '岡山県' or labels == '広島県':
            chugoku += int(numbers)
        elif labels == '徳島県' or labels == '香川県' or labels == '愛媛県' or labels == '高知県':
            sikoku += int(numbers)
        elif labels == '山口県' or labels == '福岡県' or labels == '佐賀県' or labels == '長崎県' or labels == '熊本県' or labels == '大分県' or labels == '宮崎県' or labels == '鹿児島県' or labels == '沖縄県':
            kyusyu += int(numbers)
    all_area = [kyusyu,sikoku,chugoku,kinki,tokai,center_hill,hokuriku,s_kanto,n_kanto,tohoku,hokkaido]
    # 値を再構成したので配置も再調整
    x_area = []
    for area_ranges in range(1,len(all_area) + 1):
        x_area.append(area_ranges)
    # １行１列のグラフの描画
    nrow = 1
    ncol = 1
    # col:横、row:縦
    # 選択肢の数が多い場合、選択肢数の1/3くらいでちょうど良さそう
    plt.figure(figsize=(6*ncol,round(len(all_area)/3)*nrow))
    # １つ目のsubplot領域にプロット
    plt.subplot(nrow,ncol,1)
    # 棒グラフはmatplotlibのbarメソッドを利用して作成する
    plt.barh(x_area,all_area,align='center',height=0.8)
    plt.yticks(x_area,district_label)
    # y_lim()でy軸幅を調整する（そのままだとグラフ内の上下で余白が出ることがあるため）
    # 要素数ぴったりに合わせる
    plt.ylim(0,len(all_area)+1)
    # plt.legend()：凡例を指定してしまうと謎の白箱が出てきてしまうので触れない
    # 画像の出力
    plt.savefig('./save_images/Q{}_graph.png'.format(select),dpi=200,orientation='portrait',transparent=False,pad_inches=0.0,bbox_inches='tight')

    # サブパターン：表の作成をHTMLで行う
    # jinja2テンプレート出力
    env = Environment(loader=FileSystemLoader('./', encoding='utf8'))
    tmpl = env.get_template('jinja2_templetes/templetes_extra_prefecture.tmpl')
    # 商品情報を入れるリストを宣言、初期化
    items = []
    # 回答割合データの作成（小数点一桁）
    percent = []
    for values in all_area:
        percent.append(round(values/sum(all_area)*100))

    # ループでitemsに回答ラベル、回答実数、回答割合をappendする
    for ti,tao,tat in zip(reversed(district_label),reversed(all_area),reversed(percent)):
        try:
            # isnan()は「非数の時」にFalseを返すので数字は通ってしまうので数値を扱う際は判定方法を変える
            if len(tao) == 0 or len(tat) == 0:
                pass
        except:
            items.append({'LABEL':ti,'ANSER_ONE':tao,'ANSER_TWO':tat})
    # 末尾にまとめの値を追加する
    items.append({'LABEL':'合計','ANSER_ONE':len(act_data),'ANSER_TWO':100})
    html = tmpl.render({'title':graph_title,'item':items})
    with open('jinja2_templetes/templetes_extra_prefecture.html',mode='w') as f:
        f.write(str(html))


    # seleniumでブラウザ表示（動作問題なし）
    options = Options()
    options.add_argument('--headless')
    driver = webdriver.Chrome(options=options)
    driver.get('file:///Users/t_sasaki/Documents/%E6%A5%AD%E5%8B%99/%E3%82%A2%E3%83%B3%E3%82%B1%E3%83%BC%E3%83%88/%E3%82%AB%E3%83%9F%E3%82%AA%E3%83%B3/macro_mk/jinja2_templetes/templetes_extra_prefecture.html')
    page_width = driver.execute_script('return document.body.scrollWidth')
    page_height = driver.execute_script('return document.body.scrollHeight')
    driver.set_window_size(page_width,page_height)
    driver.save_screenshot('save_images/Q{}_chart.png'.format(select))
    driver.quit()

    # 作成した画像をPPに貼り付ける
    # PPの呼び出し
    prs = Presentation('./static/pp/template.pptx')
    # 画像を中央配置にするためにスライドの幅と高さを取得する
    width = prs.slide_width
    height = prs.slide_height
    # 貼り付ける画像の読み込み
    fnms = ['./save_images/Q{}_graph.png'.format(select),'./save_images/Q{}_chart.png'.format(select)]
    # ファイルでループさせる
    for fnm in fnms:
        # スライドの追加。６は白紙のスライド
        add_s = prs.slide_layouts[6]
        # 合体？
        slide = prs.slides.add_slide(add_s)
        # 画像の挿入
        pic = slide.shapes.add_picture(fnm,0,0,width=None,height=None)
        # 画像を中心に設置
        pic.left = int((width - pic.width)/2)
        pic.top = int((height - pic.height)/2)
    # とりあえず画像サイズについては個々人で調整してもらう方向で
    prs.save('./static/pp/template.pptx')

if __name__ in '__main__':
    Camion()